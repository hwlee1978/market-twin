"""
Generate Mr.AI core capabilities deck (.pptx) for FSN strategic context.

Design: 16:9 wide · minimal editorial · navy + copper accent
Audience: FSN internal — Mr.AI as a business pillar across portfolio brands
NOT a customer pitch. NOT Le Mouton-specific.
Agent layer names: Strategist · Analyst · Synthesizer.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# ── Palette ─────────────────────────────────────────────────────
NAVY = RGBColor(0x0F, 0x1F, 0x3D)
NAVY_DEEP = RGBColor(0x08, 0x14, 0x2A)
NAVY_MID = RGBColor(0x1E, 0x34, 0x58)
INK = RGBColor(0x11, 0x19, 0x2C)
TEXT = RGBColor(0x1A, 0x22, 0x33)
TEXT_SOFT = RGBColor(0x4A, 0x54, 0x68)
TEXT_MUTE = RGBColor(0x8A, 0x93, 0xA4)
LINE = RGBColor(0xE2, 0xDD, 0xD0)
BG = RGBColor(0xFA, 0xF8, 0xF3)   # warm off-white
BG_PURE = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT = RGBColor(0xBE, 0x6E, 0x3C)  # copper
ACCENT_SOFT = RGBColor(0xF0, 0xDE, 0xC9)
EMERALD = RGBColor(0x2F, 0x6E, 0x5F)
EMERALD_SOFT = RGBColor(0xD6, 0xE8, 0xE0)

FONT_KR = "Pretendard"
FONT_EN = "Inter"
FONT_SERIF = "Cormorant Garamond"
FONT_MONO = "JetBrains Mono"

# ── Font helper ─────────────────────────────────────────────────
# python-pptx's font.name setter only writes <a:latin typeface="..">.
# When a system doesn't have that face mapped as the East-Asian (CJK)
# font, PowerPoint silently picks an entirely different fallback for
# Korean glyphs in the same run — the rendering jumps between two
# unrelated typefaces mid-line. Set ALL three slots (latin / ea / cs)
# so Pretendard is used uniformly for Latin + Korean + complex scripts.

def set_font(run, name):
    if name is None:
        return
    rPr = run._r.get_or_add_rPr()
    for slot in ('a:latin', 'a:ea', 'a:cs'):
        el = rPr.find(qn(slot))
        if el is None:
            el = etree.SubElement(rPr, qn(slot))
        el.set('typeface', name)


# Monkey-patch python-pptx so Font.name setter populates ALL 3
# typeface slots automatically — no more "I forgot to set ea" bugs
# in legacy slide-building code below.
from pptx.text.text import Font as _PptxFont
_orig_name_setter = _PptxFont.name.fset


def _patched_name_setter(self, value):
    _orig_name_setter(self, value)
    if value is None:
        return
    rPr = self._rPr
    if rPr is None:
        return
    for slot in ('a:ea', 'a:cs'):
        el = rPr.find(qn(slot))
        if el is None:
            el = etree.SubElement(rPr, qn(slot))
        el.set('typeface', value)


_PptxFont.name = property(_PptxFont.name.fget, _patched_name_setter)


# ── Glyph sanitizer ─────────────────────────────────────────────
# Pretendard ships Korean + Latin only — arrows / emoji / pictographs
# fall back to whatever's available and render as garbage characters
# on most machines. Replace them with ASCII-safe equivalents BEFORE
# any text reaches the PPTX layer.
import re as _re

_GLYPH_MAP = {
    # Arrows — Pretendard doesn't ship these, fallback fonts render
    # them as random Latin garbage on most systems. Map to ASCII.
    "→": ">", "←": "<", "↑": "+", "↓": "-",
    "⇒": ">", "⇐": "<",
    "↗": ">", "↘": ">", "↖": "<", "↙": "<",
    "•": "-",     # bullet -> hyphen
    "✓": "OK", "✗": "X", "✘": "X",
    "★": "*", "☆": "*",
    # NOTE: middle-dot (·), em-dash (—), multiplication sign (×) are
    # kept as-is. Pretendard / system Korean fonts cover them.
}
_EMOJI_RE = _re.compile(
    "[\U0001F000-\U0001FFFF"
    "☀-➿"
    "️]"
)
def _clean(txt):
    if not isinstance(txt, str):
        return txt
    for k, v in _GLYPH_MAP.items():
        txt = txt.replace(k, v)
    txt = _EMOJI_RE.sub("", txt)
    # Collapse double spaces from removed glyphs.
    txt = _re.sub(r" +", " ", txt).strip(" ")
    return txt

# ── Init 16:9 deck ──────────────────────────────────────────────
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height

BLANK = prs.slide_layouts[6]


def new_slide(bg=BG):
    s = prs.slides.add_slide(BLANK)
    fill = s.background.fill
    fill.solid()
    fill.fore_color.rgb = bg
    return s


def add_text(slide, x, y, w, h, text, *, font=FONT_KR, size=14, bold=False,
             color=TEXT, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             italic=False, letter_spacing=None):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = _clean(text)
    f = run.font
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.color.rgb = color
    set_font(run, font)
    if letter_spacing is not None:
        # set spc via xml
        rPr = run._r.get_or_add_rPr()
        rPr.set('spc', str(letter_spacing))
    return tb


def add_multi(slide, x, y, w, h, lines, *, font=FONT_KR, size=12,
              color=TEXT, anchor=MSO_ANCHOR.TOP, line_space=1.4):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.line_spacing = line_space
        text = line if isinstance(line, str) else line["text"]
        run = p.add_run()
        run.text = _clean(text)
        f = run.font
        f.size = Pt(size)
        f.color.rgb = color
        chosen_font = font
        if isinstance(line, dict):
            if line.get("bold"):
                f.bold = True
            if line.get("italic"):
                f.italic = True
            if line.get("size"):
                f.size = Pt(line["size"])
            if line.get("color"):
                f.color.rgb = line["color"]
            if line.get("font"):
                chosen_font = line["font"]
            if line.get("indent"):
                p.level = line["indent"]
        set_font(run, chosen_font)
    return tb


def add_bullet_list(slide, x, y, w, h, items, *, size=12, color=TEXT):
    """Bulleted list with em-dash marker. Avoids glyphs Pretendard
    doesn't ship (middle-dot, arrows, emoji) — those fall back to a
    random font on systems without full coverage and render as garbage.
    Also strips emoji + arrow characters defensively."""
    import re
    # Strip emoji + arrows from incoming text — keep meaning, drop glyph.
    sanitize_re = re.compile(
        r"[←-⇿⌀-⏿■-➿✀-➿"
        r"\U0001F300-\U0001FAFF️]"
    )
    lines = []
    for it in items:
        clean = sanitize_re.sub("", it).strip()
        lines.append({"text": f"—  {clean}", "size": size, "color": color})
    return add_multi(slide, x, y, w, h, lines, size=size, color=color, line_space=1.5)


def add_rect(slide, x, y, w, h, *, fill=None, line=None, line_w=0):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                Inches(x), Inches(y), Inches(w), Inches(h))
    sh.line.fill.background()
    if line is not None:
        sh.line.color.rgb = line
        sh.line.width = Pt(line_w if line_w else 0.75)
    if fill is not None:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    else:
        sh.fill.background()
    sh.shadow.inherit = False
    return sh


def add_line(slide, x1, y1, x2, y2, color=LINE, weight=0.75):
    ln = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    ln.line.color.rgb = color
    ln.line.width = Pt(weight)
    return ln


def add_pill(slide, x, y, text, *, fill=ACCENT, color=BG_PURE, size=8, w=None):
    """Small uppercase pill — section eyebrow style."""
    if w is None:
        w = max(0.7, 0.08 * len(text) + 0.4)
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                Inches(x), Inches(y), Inches(w), Inches(0.28))
    sh.line.fill.background()
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.shadow.inherit = False
    tf = sh.text_frame
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = _clean(text)
    f = r.font
    f.size = Pt(size)
    f.bold = True
    f.color.rgb = color
    set_font(r, FONT_KR)
    rPr = r._r.get_or_add_rPr()
    rPr.set('spc', '180')
    return sh


_PG_TOTAL = 38  # update after all slides built — must match prs.slides count

def page_chrome(slide, eyebrow, pg_no, total=None):
    total = total or _PG_TOTAL
    """Top eyebrow + bottom footer + thin top rule."""
    add_text(slide, 0.6, 0.42, 8, 0.3, eyebrow.upper(),
             font=FONT_KR, size=8, color=ACCENT, bold=True, letter_spacing=350)
    add_text(slide, 12.0, 0.42, 1.0, 0.3, f"{pg_no:02d} / {total:02d}",
             font=FONT_SERIF, size=10, color=TEXT_MUTE,
             italic=True, align=PP_ALIGN.RIGHT)
    add_line(slide, 0.6, 0.78, 12.73, 0.78, color=LINE, weight=0.5)
    # bottom footer
    add_text(slide, 0.6, 7.05, 8, 0.3, "Mr.AI · AI-Native CMO Infrastructure",
             font=FONT_KR, size=7.5, color=TEXT_MUTE, letter_spacing=200)
    add_text(slide, 12.0, 7.05, 1.0, 0.3, "MR.AI",
             font=FONT_SERIF, size=8, color=ACCENT, italic=True, align=PP_ALIGN.RIGHT)
    add_line(slide, 0.6, 7.02, 12.73, 7.02, color=LINE, weight=0.4)


def title_serif(slide, x, y, w, h, kr, en_em=None):
    """Big editorial-serif title with optional italic accent."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.line_spacing = 1.05
    r = p.add_run()
    r.text = _clean(kr)
    f = r.font
    f.size = Pt(34)
    f.bold = False
    f.color.rgb = NAVY
    set_font(r, FONT_KR)
    if en_em:
        r2 = p.add_run()
        r2.text = " " + _clean(en_em)
        f2 = r2.font
        f2.size = Pt(34)
        f2.italic = True
        f2.color.rgb = ACCENT
        set_font(r2, FONT_SERIF)
    return tb


def screenshot_slot(slide, x, y, w, h, label="스크린샷 슬롯"):
    """Soft placeholder with dashed border + camera glyph."""
    # outer rectangle subtle fill
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = RGBColor(0xF1, 0xEC, 0xDE)
    sh.line.color.rgb = RGBColor(0xC9, 0xB7, 0x9A)
    sh.line.width = Pt(0.75)
    sh.shadow.inherit = False
    # dashed line via xml
    spPr = sh.line._get_or_add_ln()
    prstDash = etree.SubElement(spPr, qn('a:prstDash'))
    prstDash.set('val', 'dash')
    sh.shadow.inherit = False
    tf = sh.text_frame
    tf.margin_left = Inches(0.2); tf.margin_right = Inches(0.2)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r1 = p.add_run()
    r1.text = "SCREENSHOT"
    r1.font.name = FONT_KR
    r1.font.size = Pt(9)
    r1.font.bold = True
    r1.font.color.rgb = RGBColor(0xC9, 0xB7, 0x9A)
    rPr = r1._r.get_or_add_rPr()
    rPr.set('spc', '350')
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = "\n" + _clean(label)
    r2.font.name = FONT_KR
    r2.font.size = Pt(10)
    r2.font.color.rgb = TEXT_MUTE


# ─────────────────────────────────────────────────────────────────
# Slide 1 — COVER
# ─────────────────────────────────────────────────────────────────
s = new_slide(bg=NAVY_DEEP)
# accent rule
add_rect(s, 0.6, 1.0, 0.05, 5.3, fill=ACCENT)
# eyebrow
add_text(s, 0.95, 1.0, 8, 0.35, "MR.AI",
         font=FONT_KR, size=9, color=ACCENT, bold=True, letter_spacing=600)
add_text(s, 0.95, 1.4, 8, 0.35, "AI-NATIVE CMO INFRASTRUCTURE",
         font=FONT_KR, size=9, color=RGBColor(0xC9, 0xB7, 0x9A),
         bold=True, letter_spacing=400)
# title
tb = s.shapes.add_textbox(Inches(0.95), Inches(2.4), Inches(11.5), Inches(2.8))
tf = tb.text_frame; tf.word_wrap = True
tf.margin_left = tf.margin_right = 0
tf.margin_top = tf.margin_bottom = 0
p = tf.paragraphs[0]; p.line_spacing = 1.05
r = p.add_run(); r.text = "한 명의 CMO 가\n"
r.font.name = FONT_KR; r.font.size = Pt(54); r.font.color.rgb = BG
r2 = p.add_run(); r2.text = "24시간 일한다."
r2.font.name = FONT_SERIF; r2.font.size = Pt(54)
r2.font.italic = True; r2.font.color.rgb = RGBColor(0xC9, 0xB7, 0x9A)
# subtitle
add_text(s, 0.95, 5.4, 11, 0.6,
         "FSN 포트폴리오 브랜드의 마케팅·시장조사·콘텐츠를 하나의 시스템으로.",
         font=FONT_KR, size=15, color=RGBColor(0xE7, 0xEC, 0xF4))
# bottom meta
add_line(s, 0.95, 6.6, 12.4, 6.6, color=RGBColor(0x4A, 0x5A, 0x7E), weight=0.5)
add_text(s, 0.95, 6.75, 6, 0.3, "STRATEGIC DECK · ㈜미스터에이아이",
         font=FONT_KR, size=8, color=RGBColor(0xC9, 0xB7, 0x9A),
         bold=True, letter_spacing=350)
add_text(s, 8, 6.75, 4.4, 0.3, "2026.05.29",
         font=FONT_SERIF, size=10, color=RGBColor(0xC9, 0xB7, 0x9A),
         italic=True, align=PP_ALIGN.RIGHT)

# ─────────────────────────────────────────────────────────────────
# Slide 2 — WHAT IS MR.AI
# ─────────────────────────────────────────────────────────────────
s = new_slide()
page_chrome(s, "01 · 정의", 2)
title_serif(s, 0.6, 1.05, 12.5, 1.2, "Mr.AI 는", "what it is")
add_text(s, 0.6, 2.2, 12.5, 0.5,
         "브랜드의 마케팅·시장조사·콘텐츠 운영을 하나의 메모리 위에 통합한 AI 시스템.",
         font=FONT_KR, size=15, color=TEXT_SOFT)

# 4 value props row
labels = [
    ("기억하는 AI", "Vector + Structured 메모리\n매일 컨텍스트 자동 누적"),
    ("3 계층 분업", "Strategist · Analyst · Synthesizer\n역할별 책임 명확"),
    ("자동 운영", "크롤 · 브리핑 · 콘텐츠 생성\n발행까지 매일 자동"),
    ("측정 가능", "KPI 피드백 루프\n모든 결과가 다음 학습 input"),
]
top = 3.5; left = 0.6; box_w = 3.0; gap = 0.18
for i, (h, body) in enumerate(labels):
    x = left + i * (box_w + gap)
    add_rect(s, x, top, box_w, 2.7, fill=BG_PURE, line=LINE, line_w=0.5)
    add_pill(s, x + 0.25, top + 0.3, f"0{i+1}",
             fill=ACCENT_SOFT, color=NAVY, size=7, w=0.5)
    add_text(s, x + 0.25, top + 0.75, box_w - 0.5, 0.5, h,
             font=FONT_KR, size=15, color=NAVY, bold=True)
    add_text(s, x + 0.25, top + 1.4, box_w - 0.5, 1.2, body,
             font=FONT_KR, size=11, color=TEXT_SOFT)


# ─────────────────────────────────────────────────────────────────
# Slide 3 — ARCHITECTURE (Strategist · Analyst · Synthesizer)
# ─────────────────────────────────────────────────────────────────
s = new_slide()
page_chrome(s, "02 · 아키텍처 · 3 계층", 3)
title_serif(s, 0.6, 1.05, 12.5, 1.2, "3 계층 에이전트", "architecture")
add_text(s, 0.6, 2.2, 12.5, 0.5,
         "한 AI 가 다 하면 — 책임도 학습도 흐려진다. Mr.AI 는 의도적으로 셋으로 나뉘어 있다.",
         font=FONT_KR, size=13, color=TEXT_SOFT)

layers = [
    ("STRATEGIST", "전략", "분기 KPI · 우선순위 · 시장 진출 결정.\nMarket Twin 결과를 의사결정으로 변환.",
     NAVY, BG, RGBColor(0xC9, 0xB7, 0x9A)),
    ("ANALYST", "분석", "채널·캠페인·세그먼트 단위 기획.\n페르소나 데이터를 콘텐츠 spec 으로 변환.",
     ACCENT_SOFT, NAVY, ACCENT),
    ("SYNTHESIZER", "합성", "실제 콘텐츠·이미지·시뮬·발행.\n결과를 KPI 피드백 루프로 다시 보고.",
     BG_PURE, NAVY, TEXT_SOFT),
]
top = 3.4; left = 0.6; box_w = 4.0; gap = 0.2
for i, (en, kr, body, bg, fg, accent) in enumerate(layers):
    x = left + i * (box_w + gap)
    add_rect(s, x, top, box_w, 3.0, fill=bg,
             line=LINE if bg == BG_PURE else bg, line_w=0.5)
    add_text(s, x + 0.4, top + 0.35, box_w - 0.8, 0.3, f"LAYER 0{i+1}",
             font=FONT_KR, size=8, color=accent, bold=True, letter_spacing=350)
    # layer kr/en
    tb = s.shapes.add_textbox(Inches(x + 0.4), Inches(top + 0.7),
                              Inches(box_w - 0.8), Inches(1.0))
    tf = tb.text_frame; tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]; p.line_spacing = 1.05
    r = p.add_run(); r.text = en
    r.font.name = FONT_SERIF; r.font.size = Pt(28)
    r.font.italic = True; r.font.color.rgb = fg
    p2 = tf.add_paragraph()
    r2 = p2.add_run(); r2.text = kr
    r2.font.name = FONT_KR; r2.font.size = Pt(14); r2.font.color.rgb = fg
    # body
    add_text(s, x + 0.4, top + 1.95, box_w - 0.8, 1.0, body,
             font=FONT_KR, size=10.5, color=fg)


# ─────────────────────────────────────────────────────────────────
# Slide 4 — MEMORY (Vector + Structured)
# ─────────────────────────────────────────────────────────────────
s = new_slide()
page_chrome(s, "03 · 메모리 시스템", 4)
title_serif(s, 0.6, 1.05, 12.5, 1.2, "어제를 기억한다", "memory")
add_text(s, 0.6, 2.2, 12.5, 0.5,
         "Vector 임베딩 + Structured 사실 · 두 종류로 누적. 모든 답변이 워크스페이스 컨텍스트 안에서 나옴.",
         font=FONT_KR, size=13, color=TEXT_SOFT)

# Two columns
cols = [
    ("STRUCTURED", "검증된 사실", "PostgreSQL 정형 테이블",
     ["타입 · 제목 · 본문 · 태그 4 필드",
      "임원 결정 · KPI · 가격 · 결정사항",
      "PDF 자동 추출기로 시드 (50p → 메모리 30-60건)"]),
    ("VECTOR (pgvector)", "의미 검색", "1,536차원 임베딩 · cosine similarity",
     ["자유 질의 → 의미적 유사 메모리 자동 추출",
      "지난 대화 · 시뮬 결과 모두 검색 가능",
      "워크스페이스 단위 영구 누적"]),
]
top = 3.4
for i, (en, kr, sub, items) in enumerate(cols):
    x = 0.6 + i * 6.25
    add_rect(s, x, top, 6.05, 3.0, fill=BG_PURE, line=LINE, line_w=0.5)
    add_text(s, x + 0.4, top + 0.35, 5.25, 0.3, en,
             font=FONT_KR, size=9, color=ACCENT, bold=True, letter_spacing=350)
    add_text(s, x + 0.4, top + 0.7, 5.25, 0.5, kr,
             font=FONT_KR, size=19, color=NAVY, bold=True)
    add_text(s, x + 0.4, top + 1.2, 5.25, 0.3, sub,
             font=FONT_KR, size=10.5, color=TEXT_SOFT, italic=True)
    add_bullet_list(s, x + 0.4, top + 1.7, 5.25, 1.4, items, size=11)


# ─────────────────────────────────────────────────────────────────
# Slide 5 — WORKSPACE MODEL
# ─────────────────────────────────────────────────────────────────
s = new_slide()
page_chrome(s, "04 · 워크스페이스 모델", 5)
title_serif(s, 0.6, 1.05, 12.5, 1.2, "브랜드 단위 격리", "workspace")
add_text(s, 0.6, 2.2, 12.5, 0.6,
         "FSN 포트폴리오의 각 브랜드 = 워크스페이스 1개. 데이터·메모리·채널·KPI 모두 워크스페이스 단위로 분리.",
         font=FONT_KR, size=13, color=TEXT_SOFT)

# Visual: 1 account → N workspaces
add_rect(s, 0.6, 3.2, 2.6, 3.3, fill=NAVY, line=NAVY)
add_text(s, 0.85, 3.45, 2.1, 0.3, "FSN ACCOUNT",
         font=FONT_KR, size=8, color=RGBColor(0xC9, 0xB7, 0x9A),
         bold=True, letter_spacing=350)
add_text(s, 0.85, 3.85, 2.1, 0.5, "본사 계정",
         font=FONT_KR, size=18, color=BG, bold=True)
add_text(s, 0.85, 4.5, 2.1, 2, "Super admin 권한.\n포트폴리오 전체 가시성.\nMember 초대 · 권한 관리.",
         font=FONT_KR, size=10.5, color=RGBColor(0xC9, 0xB7, 0x9A))

# Arrow
add_line(s, 3.4, 4.85, 4.3, 4.85, color=ACCENT, weight=1.5)
add_line(s, 4.1, 4.7, 4.3, 4.85, color=ACCENT, weight=1.5)
add_line(s, 4.1, 5.0, 4.3, 4.85, color=ACCENT, weight=1.5)

# Workspaces
ws_list = [
    ("브랜드 A", "Footwear · KR + JP"),
    ("브랜드 B", "Apparel · KR + US"),
    ("브랜드 C", "Cosmetics · 글로벌"),
    ("브랜드 D · 신규", "검토 단계"),
]
top = 3.2; left = 4.5
for i, (n, sub) in enumerate(ws_list):
    y = top + i * 0.85
    add_rect(s, left, y, 8.3, 0.75, fill=BG_PURE, line=LINE)
    add_text(s, left + 0.3, y + 0.13, 3.5, 0.25, "WORKSPACE",
             font=FONT_KR, size=7, color=ACCENT, bold=True, letter_spacing=300)
    add_text(s, left + 0.3, y + 0.38, 4, 0.35, n,
             font=FONT_KR, size=13, color=NAVY, bold=True)
    add_text(s, left + 4.5, y + 0.25, 3.5, 0.35, sub,
             font=FONT_KR, size=11, color=TEXT_SOFT, align=PP_ALIGN.RIGHT)

# Bottom note
add_text(s, 0.6, 6.6, 12.5, 0.3,
         "데이터 격리 · Row-Level Security · 페르소나 풀만 글로벌 공유 (cold start 비용 절감)",
         font=FONT_KR, size=10.5, color=TEXT_MUTE, italic=True)

# ─────────────────────────────────────────────────────────────────
# Build a section divider helper
# ─────────────────────────────────────────────────────────────────

def section_divider(part_no, kr, en_em):
    next_pg()  # increment global counter so feature slides stay correct
    s = new_slide(bg=NAVY)
    # accent rule
    add_rect(s, 0.6, 2.5, 0.05, 2.3, fill=ACCENT)
    add_text(s, 0.95, 2.5, 6, 0.35, f"PART {part_no:02d}",
             font=FONT_KR, size=9, color=ACCENT, bold=True, letter_spacing=400)
    tb = s.shapes.add_textbox(Inches(0.95), Inches(2.95), Inches(12), Inches(1.5))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    p = tf.paragraphs[0]; p.line_spacing = 1.05
    r = p.add_run(); r.text = _clean(kr) + " "
    r.font.name = FONT_KR; r.font.size = Pt(40); r.font.color.rgb = BG
    r2 = p.add_run(); r2.text = _clean(en_em)
    r2.font.name = FONT_SERIF; r2.font.size = Pt(40)
    r2.font.italic = True
    r2.font.color.rgb = RGBColor(0xC9, 0xB7, 0x9A)


# ── Feature slide layout helper ──────────────────────────────────
# Global page counter — cover is page 1; every feature/section slide
# increments. Avoids manual renumbering when slides are inserted.
_pg = [1]  # mutable so closures can update

def next_pg():
    _pg[0] += 1
    return _pg[0]

def feature_slide(*args, ss_w=5.5):
    """Standard feature slide: title left + screenshot slot right.
    Tolerates a legacy leading int arg (old call sites)."""
    if args and isinstance(args[0], int):
        args = args[1:]
    part, kr_title, en_em, timing, summary, bullets, ss_label = args
    pg = next_pg()
    s = new_slide()
    page_chrome(s, part, pg)
    # title
    title_serif(s, 0.6, 1.0, 7.3, 1.2, kr_title, en_em)
    # timing pill
    if timing:
        add_pill(s, 0.6, 2.25, timing, fill=ACCENT_SOFT, color=NAVY,
                 size=8, w=max(1.0, 0.075 * len(timing) + 0.25))
    # summary
    add_text(s, 0.6, 2.75, 7.3, 0.7, summary,
             font=FONT_KR, size=13.5, color=TEXT_SOFT)
    # bullets
    add_bullet_list(s, 0.6, 3.85, 7.3, 2.7, bullets, size=12)
    # right side screenshot slot
    screenshot_slot(s, 8.2, 1.05, 4.7, 5.5, ss_label)


# ─────────────────────────────────────────────────────────────────
# Slide 6 — Login & workspace switching
# ─────────────────────────────────────────────────────────────────
feature_slide(
    6, "00 · 진입",
    "로그인 + 워크스페이스 전환", "entry",
    "1초 (전환)",
    "한 계정으로 N 브랜드 워크스페이스 운영. 좌측 상단 스위처 한 번에 컨텍스트 전환.",
    [
        "Google OAuth · 이메일·비밀번호 2 경로",
        "워크스페이스 스위처 — 모든 페이지 동일 위치",
        "전환 시 모달 안내 (메모리 · 시뮬 · 채널 재로딩 표시)",
        "ID 기반 경로는 새 워크스페이스의 list 페이지로 자동 redirect",
    ],
    "/dashboard · 좌측 상단 워크스페이스 드롭다운 + 전환 오버레이",
)


# ─────────────────────────────────────────────────────────────────
# Part II Setup phase divider
# ─────────────────────────────────────────────────────────────────
section_divider(2, "셋업 단계 ·", "first hour")

# ─────────────────────────────────────────────────────────────────
# Slide 7 — Auto-Seed
# ─────────────────────────────────────────────────────────────────
feature_slide(
    7, "05 · Auto-Seed",
    "Auto-Seed", "1-minute genesis",
    "45초 ~ 2분 · $0.10-0.30 / 회",
    "회사명 한 줄 입력 → 웹 리서치 + LLM 합성 → 8단계 컨텍스트 자동 생성.",
    [
        "8 답변: business · scale · products · channels · competitors · executive · decisions · kpi",
        "각 답변 100-600자, 검색 결과의 브랜드명·숫자·가격 보존",
        "needsReview 자동 태그 — 임원 검토 사항 명시",
        "메모리 시드 + Mr.AI 모든 후속 답변의 컨텍스트로 즉시 사용",
    ],
    "/mr-ai · Auto-Seed 모달 + 8 답변 결과 카드",
)

# ─────────────────────────────────────────────────────────────────
# Slide — Mr.AI Chat (central conversation surface)
# ─────────────────────────────────────────────────────────────────
feature_slide(
    "06 · 채팅",
    "Mr.AI 채팅", "conversation",
    "상시 · 무제한",
    "워크스페이스 컨텍스트 + 메모리 + 지식 그래프 위의 대화. 일반 챗봇과 다름.",
    [
        "어제·이번 주·분기 누적 메모리를 자동 호출하여 답변",
        "PDF 드래그&드롭으로 메모리 즉시 시드",
        "코드·표·차트·이미지 모든 응답 형식 지원",
        "Strategist 계층 — 분기 전략/우선순위 직접 토론 가능",
    ],
    "/mr-ai · 채팅 화면 + 메모리 컨텍스트 사이드패널",
)


# ─────────────────────────────────────────────────────────────────
# Slide — Memory seeding (PDF extractor)
# ─────────────────────────────────────────────────────────────────
feature_slide(
    "07 · PDF 메모리 추출",
    "메모리 시딩", "memory seeding",
    "15-30초 / PDF · $0.10",
    "사업계획서 · IR 자료 · 컨설팅 보고서 PDF 드래그&드롭 → '검증 가능한 사실' 만 자동 추출.",
    [
        "50p PDF → 메모리 30-60건 자동 등재",
        "Claude Vision + 텍스트 추출 · 표 · 차트 캡션 포함",
        "type 자동 분해: fact / context / decision / preference",
        "추출된 사실은 모든 후속 콘텐츠 · SEO · 시뮬의 컨텍스트",
    ],
    "/mr-ai · 채팅에 PDF 드래그 후 추출 진행 화면",
)

# ─────────────────────────────────────────────────────────────────
# Slide 9 — Brand asset library
# ─────────────────────────────────────────────────────────────────
feature_slide(
    9, "07 · 브랜드 자산 라이브러리",
    "브랜드 자산", "brand assets",
    "10분 (업로드) · 자동 분석 60초",
    "제품 사진 · 로고 · 룩북 · 앰배서더 · 패턴 5 유형으로 자동 정리. 모든 콘텐츠 생성의 시각 reference.",
    [
        "Claude Vision 으로 각도 · 색상 · 로고 위치 자동 메타데이터",
        "제품 프로필 (다음 슬라이드) 의 입력으로 직결",
        "이미지 생성 시 reference 사진 자동 첨부 → 무브랜드 강제 + 후처리 로고 합성",
        "워크스페이스별 격리 · 영구 캐시",
    ],
    "/mr-ai/brand · 5 유형 그리드 + 업로드 모달",
)

# ─────────────────────────────────────────────────────────────────
# Slide 10 — Product profile (Vision)
# ─────────────────────────────────────────────────────────────────
feature_slide(
    10, "08 · 제품 프로필",
    "제품 프로필", "vision spec",
    "$0.02-0.05 · 15-30초",
    "Claude Vision 이 제품 사진 2-5장을 합쳐 단일 정밀 spec 으로 변환. 모든 이미지 생성의 정확한 blueprint.",
    [
        "category · description · 실루엣 · 소재 · 색상 (hex) · 로고 위치 자동 추출",
        "물리 제품 없는 워크스페이스 (SaaS·디지털) 는 manual 카테고리 선택",
        "이미지 생성 시 시각 spec 자동 prepend → 다른 브랜드 실루엣으로 변형 차단",
        "가짜 브랜드 글자 환각 (가장 흔한 실패 모드) 완전 차단",
    ],
    "/mr-ai/brand · 제품 프로필 카드 (추출 결과 + 컬러 스와치)",
)


# ─────────────────────────────────────────────────────────────────
# Part III Daily operation divider
# ─────────────────────────────────────────────────────────────────
section_divider(3, "매일 운영 ·", "daily operation")


# ─────────────────────────────────────────────────────────────────
# Slide 12 — Auto-crawl
# ─────────────────────────────────────────────────────────────────
feature_slide(
    12, "09 · 자동 크롤",
    "자동 크롤 소스", "daily crawl",
    "매일 02:30 KST · 자동",
    "자사 사이트 · 자사 브랜드 뉴스 · 경쟁사 페이지 · 카테고리 트렌드 자동 수집. 변동만 새 메모리로 저장.",
    [
        "카테고리별 프리셋 — footwear / cosmetics / saas_digital 등 자동 노출",
        "Google News RSS · 직접 페이지 fetch 둘 다 지원 (Cloudflare 우회 포함)",
        "프리셋 다중 선택 → 일괄 등록 (그룹 전체 선택 가능)",
        "다음 날 Daily Briefing 에 자동 반영 — 경쟁사 신상 출시 알림 1줄",
    ],
    "/mr-ai/brand · 크롤 프리셋 일괄 선택 모달 + 활성 소스 리스트",
)

# ─────────────────────────────────────────────────────────────────
# Slide 13 — Daily Briefing
# ─────────────────────────────────────────────────────────────────
feature_slide(
    13, "10 · Daily Briefing",
    "Daily Briefing", "08:00 KST",
    "매일 08:00 KST · Slack / Email",
    "어제 요약 · 오늘 챙길 것 · 주의 신호 3 섹션. 메모리 + 크롤 + KPI 자동 합성.",
    [
        "출처 명확 — 추측 일반론 없음. '정보 부족' 명시",
        "Slack OAuth 1회 연결 → 매일 자동 발송",
        "워크스페이스 멤버 전체 또는 특정 채널 발송",
        "임원의 매일 아침 45분 → 3분으로 단축 (연 250시간)",
    ],
    "Slack 채널의 실제 브리핑 메시지 예시",
)

# ─────────────────────────────────────────────────────────────────
# Slide 14 — SEO + LLM Visibility
# ─────────────────────────────────────────────────────────────────
feature_slide(
    14, "11 · SEO + LLM 가시성",
    "검색 가시성", "search · llm",
    "30분 셋업 · 월간 audit",
    "Google + 네이버 + ChatGPT · Claude · Gemini 모두에서 브랜드 인용도 측정. 2026 의 새 KPI.",
    [
        "GSC · GA4 · 네이버 서치어드바이저 단계별 가이드 (도움말 펼침)",
        "LLM Visibility Audit: 3 LLM × 20 쿼리 = 60건 병렬 실행 ($0.30)",
        "알리아스 인식 (한글·영문·합성형 통합 카운트)",
        "경쟁사 비교 score + 시간 추세 (월간 자동 측정)",
    ],
    "/mr-ai/analytics · LLM Visibility 점수 + 경쟁사 비교 차트",
)

# ─────────────────────────────────────────────────────────────────
# Slide 15 — Channels
# ─────────────────────────────────────────────────────────────────
feature_slide(
    15, "12 · 마케팅 채널",
    "마케팅 채널 셋업", "channels",
    "10분 (Auto-Seed) + 채널별 5분 보정",
    "X · IG · YouTube · 네이버 블로그 · TikTok 5 채널 자동 시드. 각자 다른 페르소나 · 다른 포스팅 스타일.",
    [
        "워크스페이스 이름 기반 핸들 자동 생성 (@brand_global · @brand.kr 등)",
        "채널별 한 줄 페르소나 + posting_style spec",
        "콘텐츠 생성 시 채널 spec 자동 prepend → 인스타용이 X 에 그대로 안 가게",
        "각 채널 카드에 '이 공간 페르소나' 미리보기",
    ],
    "/mr-ai/channels · 5 채널 카드 그리드 + 페르소나 미리보기",
)


# ─────────────────────────────────────────────────────────────────
# Slide — Members + roles
# ─────────────────────────────────────────────────────────────────
feature_slide(
    "13 · 멤버 · 권한",
    "멤버 초대 + 권한", "team",
    "1회 셋업 · 분 단위",
    "워크스페이스 단위 4 역할. Owner · Admin · Analyst · Viewer.",
    [
        "Owner: 모든 권한 + 결제 · 워크스페이스 삭제",
        "Admin: 멤버 초대 · 채널 · 통합 · 콘텐츠 발행",
        "Analyst: 시뮬 실행 · 콘텐츠 작성 (발행 제한 옵션)",
        "Viewer: 결과 조회 전용 — 외부 파트너 공유용",
    ],
    "/settings/team · 멤버 리스트 + 역할 드롭다운",
)


# ─────────────────────────────────────────────────────────────────
# Part IV Content phase divider
# ─────────────────────────────────────────────────────────────────
section_divider(4, "콘텐츠 생산 ·", "content engine")

# ─────────────────────────────────────────────────────────────────
# Slide 17 — Content drafter A/B/C × 5 formats
# ─────────────────────────────────────────────────────────────────
feature_slide(
    17, "13 · 콘텐츠 드래프터",
    "AI 콘텐츠 생성", "drafter",
    "30-60초 · $0.05/변형",
    "주제 한 줄 + 캠페인 라벨 → A/B/C 3 변형 자동 생성. 5 포맷 (기본·비교·Q&A·정의·리스트) 중 선택.",
    [
        "각 변형: body · hashtags · CTA · image_prompt · SEO 메타 동시 생성",
        "Bilingual: target market 언어 + 한국어 번역 (오퍼레이터 확인용)",
        "SEO 점수 + LLM-SEO 점수 자동 계산 (인용 가능성)",
        "워크스페이스 톤 학습 → 매주 winner 패턴 spec 자동 반영",
    ],
    "/mr-ai/channels/[id] · 생성 모달 + 결과 카드 3개 (SEO 배지 포함)",
)

# ─────────────────────────────────────────────────────────────────
# Slide 18 — Image generation
# ─────────────────────────────────────────────────────────────────
feature_slide(
    18, "14 · 이미지 자동 생성",
    "이미지 생성", "gpt-image-1",
    "15-30초 / 이미지 · $0.042",
    "콘텐츠 카드 한 클릭 → 카테고리별 시각 spec 자동 적용 → gpt-image-1 → 후처리 로고 합성.",
    [
        "카테고리별 다른 spec: footwear · cosmetics · electronics · food 등",
        "제품 무브랜드 강제 + 가짜 글자 환각 차단",
        "실루엣 · 색상 · 로고 위치 reference 사진 정확 매칭",
        "1 프레임 / 캐러셀 모두 지원 · 플랫폼별 비율 (16:9 · 1:1 · 9:16)",
    ],
    "콘텐츠 카드에 생성된 이미지 + image_prompt 영문/한글",
)

# ─────────────────────────────────────────────────────────────────
# Slide — Content edit / unpublish
# ─────────────────────────────────────────────────────────────────
feature_slide(
    "17 · 편집 · 발행 취소",
    "콘텐츠 편집 및 라이프사이클", "edit revision",
    "즉시",
    "AI 가 만든 초안은 출발점. 임원이 한 줄 수정해도 — 모든 변경이 추적되고 KPI 루프에 반영.",
    [
        "✏️ 편집 — body · CTA · 해시태그 · 이미지 프롬프트 · SEO 메타 모두 인라인",
        "이미지 프롬프트 수정 후 재생성 → 새 컷 (기존 보존)",
        "발행 confirm 다이얼로그 — body 80자 미리보기로 실수 클릭 방지",
        "발행 취소 — 누적 좋아요 · 댓글 · 시뮬 기록 함께 제거",
    ],
    "콘텐츠 카드 ✏️ 편집 모달 + 발행 취소 링크",
)


# ─────────────────────────────────────────────────────────────────
# Slide — Persona reaction simulation
# ─────────────────────────────────────────────────────────────────
feature_slide(
    "18 · 페르소나 반응 시뮬",
    "발행 전 검증", "persona sim",
    "30명: 60초 · $0.10 | 100명: 3분 · $0.36",
    "A/B/C 중 어느 변형이 가장 잘 반응하는지 — 발행 전에 페르소나 풀에서 사전 검증.",
    [
        "채널 target market 매칭 페르소나 자동 샘플링",
        "like · comment · share · save · scroll-past 5 신호 집계",
        "페르소나별 voice (1-2 문장) 수집 — 거부 / 신뢰 이유 명시",
        "winner 자동 마킹 → KPI 루프 학습 input",
    ],
    "콘텐츠 카드 시뮬 결과 표 (A/B/C × like_rate × voice 5개)",
)

# ─────────────────────────────────────────────────────────────────
# Slide 20 — Publish + KPI loop
# ─────────────────────────────────────────────────────────────────
feature_slide(
    20, "16 · 발행 KPI 루프",
    "발행 및 학습", "publish learn",
    "즉시 · 매주 자동 학습",
    "Winner 변형 발행 → 누적 좋아요 · 댓글 · 공유 자동 추적 → 다음 콘텐츠 spec 에 자동 반영.",
    [
        "발행 시 confirm 다이얼로그 (실수 클릭 방지) + 발행 취소 가능",
        "가상 피드 시뮬 cron — follower 성장 · 반응 누적",
        "4주 후 톤 정량 규칙 도출 — 외주 brief 1 page 로 추출 가능",
        "v0.2: X · LinkedIn · IG OAuth 실 발행 연동",
    ],
    "콘텐츠 카드 발행 배지 + 누적 KPI · 매주 winner 패턴 화면",
)


# ─────────────────────────────────────────────────────────────────
# Part V Market Twin section divider
# ─────────────────────────────────────────────────────────────────
section_divider(5, "Market Twin ·", "market validation")


# ─────────────────────────────────────────────────────────────────
# Slide — Market Twin: Project wizard
# ─────────────────────────────────────────────────────────────────
feature_slide(
    "21 · 프로젝트 생성",
    "Market Twin · 프로젝트", "wizard",
    "1분 입력",
    "제품·국가·가격·경쟁사 URL — 6 필드 한 화면. 경쟁사 가격은 자동 추출.",
    [
        "출시 모국 + 후보 진출국 (도메스틱 포함 가능, Top-2 추천에서는 제외)",
        "경쟁사 URL → Tavily + LLM 으로 가격 자동 추출 → anchoring 데이터",
        "카테고리별 규제 자동 검토 (footwear · 식품 · 화장품 등)",
        "사진/룩북 첨부 — 시각 reference 로 페르소나 시뮬에 반영",
    ],
    "/projects/new · 6 단계 위자드 + 진행률",
)


# ─────────────────────────────────────────────────────────────────
# Slide — Market Twin: Tier selection
# ─────────────────────────────────────────────────────────────────
feature_slide(
    "22 · 티어 선택",
    "시뮬 정밀도", "4 tiers",
    "Hypothesis ~ Deep · ₩2k ~ ₩50k",
    "의사결정 단계에 맞춘 4 티어. 페르소나 수 · 시뮬 수 · 시간 · 비용 모두 다름.",
    [
        "Hypothesis (200명 · 5-8분 · ₩2k) — 방향 빠르게",
        "Decision (1,200명 · 10-15분 · ₩8k) — 임원회의 사전 자료",
        "Decision+ (3,000명 · 15-22분 · ₩20k) — 본격 진출 결정",
        "Deep (5,000명 · 22-30분 · ₩50k) — 이사회 · IR 자료",
    ],
    "RunEnsemble 버튼 모달 · 4 티어 카드 비교",
)


# ─────────────────────────────────────────────────────────────────
# Slide — Market Twin: Multi-LLM ensemble
# ─────────────────────────────────────────────────────────────────
feature_slide(
    "23 · Multi-LLM 앙상블",
    "한 LLM 의 편향을 가둠", "ensemble",
    "병렬 실행 · 합의 등급 자동",
    "모든 시뮬을 3 LLM (Anthropic · OpenAI · DeepSeek) round-robin 으로 실행 → 합의 등급 자동 표시.",
    [
        "STRONG / WEAK / NO CONSENSUS 3 등급 — 그대로 의사결정 신뢰도",
        "DeepSeek 의 아시아 시장 인식 강점 자동 가중",
        "narrative merge LLM 이 N sim 결과를 단일 합의 narrative 로 통합",
        "Cloud Run worker · Vercel 병렬로 60초 제약 우회",
    ],
    "결과 페이지 · 합의 등급 배지 + per-LLM 점수 비교",
)


# ─────────────────────────────────────────────────────────────────
# Slide — Market Twin: Results dashboard
# ─────────────────────────────────────────────────────────────────
feature_slide(
    "24 · 결과 대시보드",
    "추천 진출국 + 가격", "results",
    "즉시 (시뮬 완료 시)",
    "Top-2 추천국 + 국가별 5 컴포넌트 점수 + 가격 곡선 + 페르소나 voice — 한 페이지 요약.",
    [
        "5 컴포넌트: 가격 fit · 채널 fit · 페르소나 fit · 경쟁 강도 · 규제 friction",
        "5 지점 가격 곡선 — sweet spot 자동 표시 (전환 × 매출 최대화)",
        "거부 / 신뢰 voice 자동 카운트 + 직접 인용 표시",
        "10+ 탭 — Overview · 국가 · 가격 · 페르소나 · 의사결정 보조 · 리스크 · 액션 · 데이터 등",
    ],
    "/projects/[id]/results · 첫 화면 추천 카드 + 점수표",
)


# ─────────────────────────────────────────────────────────────────
# Slide — Market Twin: Persona insights + chat
# ─────────────────────────────────────────────────────────────────
feature_slide(
    "25 · 페르소나 인터뷰",
    "1,200명 voice", "persona chat",
    "즉시 + 추가 인터뷰 무료",
    "각 페르소나 카드 → 💬 채팅으로 정성 인터뷰 모드. 페르소나 정체성 유지한 추가 질의.",
    [
        "국가 · 직군 · 연령 · 의향 · 가격 민감도 필터링",
        "거부 / 신뢰 voice — 자유 형식 1-2문장 자동 수집",
        "💬 채팅 — 그 페르소나의 정체성으로 추가 질문 (인터뷰처럼)",
        "필터링한 페르소나 그룹 전체에 같은 질문 일괄 가능",
    ],
    "페르소나 탭 · 필터 + 카드 + 채팅 모달",
)


# ─────────────────────────────────────────────────────────────────
# Slide — Market Twin: PDF reports
# ─────────────────────────────────────────────────────────────────
feature_slide(
    "26 · PDF 리포트",
    "임원용 · 상세 · 교차검증", "3 variants",
    "30-90초 / 종 · 무료",
    "한 클릭 → 3 종 PDF 자동 생성. 폰트 Pretendard · 미니멀 editorial 디자인.",
    [
        "임원용 (5-7p): 1 페이지 요약 + 권장사항 + 가격 곡선 + 위험",
        "상세 (15-25p): 모든 국가 + 페르소나 voice 30건+ + 시나리오 분석",
        "교차검증 (10-15p): 방법론 + Multi-LLM 합의 등급 + 외부 데이터 출처",
        "한국어 · 영문 · 시장 언어 (JP · CN · TW) 모두 지원",
    ],
    "결과 페이지 PDF 메뉴 + 다운로드된 임원용 PDF 첫 페이지",
)


# ─────────────────────────────────────────────────────────────────
# Part VI Integration & value divider
# ─────────────────────────────────────────────────────────────────
section_divider(6, "통합 · 200% 활용 ·", "value")

# ─────────────────────────────────────────────────────────────────
# Slide 22 — External integrations
# ─────────────────────────────────────────────────────────────────
s = new_slide()
page_chrome(s, "17 · 외부 통합", next_pg())
title_serif(s, 0.6, 1.05, 12.5, 1.2, "외부 시스템 통합", "integrations")
add_text(s, 0.6, 2.2, 12.5, 0.5,
         "워크스페이스 안의 데이터 · 결정 · KPI 가 외부 시스템과 양방향 동기화.",
         font=FONT_KR, size=13, color=TEXT_SOFT)

# Integration grid 2x3
integs = [
    ("HubSpot", "CRM · Deals", "거래 · 컨택 · 파이프라인\n자동 동기화 (OAuth)"),
    ("Slack", "Briefing · Alert", "Daily Briefing 발송\n채널별 알림 라우팅"),
    ("Resend", "Email", "결과 알림 · 위클리 리포트\n자동 발송"),
    ("Google Search Console", "SEO Indexing", "콘텐츠 인덱싱 상태\n자동 추적"),
    ("Google Analytics 4", "Behavior · Vitals", "방문자 행동 · Web Vitals\nKPI 루프 입력"),
    ("네이버 서치어드바이저", "KR Search", "한국 검색 인덱싱\nSitemap 자동 ping"),
]
top = 3.3; box_w = 4.0; box_h = 1.85; gap = 0.18
for i, (n, sub, body) in enumerate(integs):
    col = i % 3; row = i // 3
    x = 0.6 + col * (box_w + gap)
    y = top + row * (box_h + 0.18)
    add_rect(s, x, y, box_w, box_h, fill=BG_PURE, line=LINE, line_w=0.5)
    add_text(s, x + 0.3, y + 0.2, box_w - 0.6, 0.25, sub.upper(),
             font=FONT_KR, size=8, color=ACCENT, bold=True, letter_spacing=300)
    add_text(s, x + 0.3, y + 0.55, box_w - 0.6, 0.4, n,
             font=FONT_KR, size=15.5, color=NAVY, bold=True)
    add_text(s, x + 0.3, y + 1.05, box_w - 0.6, 0.7, body,
             font=FONT_KR, size=10.5, color=TEXT_SOFT)


# ─────────────────────────────────────────────────────────────────
# Slide 23 — Market Twin integration
# ─────────────────────────────────────────────────────────────────
s = new_slide()
page_chrome(s, "18 · Market Twin 연결", next_pg())
title_serif(s, 0.6, 1.05, 12.5, 1.2, "Market Twin 과 한 워크스페이스", "market twin")
add_text(s, 0.6, 2.2, 12.5, 0.7,
         "시장 진출 검증 (Market Twin) 의 결과가 → 메모리 → Mr.AI 콘텐츠 생성의 컨텍스트로 즉시 흐름.\n"
         "1 주: 시뮬 / 2 주: 결과 메모리 등재 / 3 주: 콘텐츠 생성 / 4 주: 발행 및 KPI 측정.",
         font=FONT_KR, size=13, color=TEXT_SOFT)

# 4-week timeline
weeks = [
    ("WEEK 1", "Market Twin 시뮬", "Decision+ 티어 · 3,000 페르소나\n24개국 평가 · 추천 진출국 + 가격"),
    ("WEEK 2", "메모리 자동 등재", "결과를 워크스페이스 메모리에\ntype=decision / fact / context"),
    ("WEEK 3", "콘텐츠 자동 생성", "거부·신뢰 voice 컨텍스트로\n타겟 시장 언어 12+ 변형"),
    ("WEEK 4", "발행 및 KPI 측정", "winner 시뮬 검증 후 발행\n실제 KPI vs 시뮬 예측 비교"),
]
top = 3.6; box_w = 3.0; gap = 0.18
for i, (w, t, body) in enumerate(weeks):
    x = 0.6 + i * (box_w + gap)
    add_rect(s, x, top, box_w, 2.7, fill=BG_PURE, line=LINE, line_w=0.5)
    add_pill(s, x + 0.25, top + 0.3, w, fill=NAVY, color=BG, size=8,
             w=1.0)
    add_text(s, x + 0.25, top + 0.85, box_w - 0.5, 0.6, t,
             font=FONT_KR, size=14, color=NAVY, bold=True)
    add_text(s, x + 0.25, top + 1.55, box_w - 0.5, 1.0, body,
             font=FONT_KR, size=10.5, color=TEXT_SOFT)
    # arrow
    if i < len(weeks) - 1:
        ax = x + box_w + 0.02
        add_line(s, ax, top + 1.35, ax + 0.13, top + 1.35,
                 color=ACCENT, weight=1.5)
        add_line(s, ax + 0.08, top + 1.25, ax + 0.13, top + 1.35,
                 color=ACCENT, weight=1.5)
        add_line(s, ax + 0.08, top + 1.45, ax + 0.13, top + 1.35,
                 color=ACCENT, weight=1.5)

# 24-month learning note
add_text(s, 0.6, 6.6, 12.5, 0.3,
         "24개월 후 — 시뮬 노이즈 ±15% → ±5%. 자사 실 진출 데이터로 자체 보정되는 의사결정 인프라.",
         font=FONT_KR, size=10.5, color=TEXT_MUTE, italic=True)


# ─────────────────────────────────────────────────────────────────
# Slide 24 — 200% utilization (one-week scenario)
# ─────────────────────────────────────────────────────────────────
s = new_slide()
page_chrome(s, "19 · 200% 활용 시나리오", next_pg())
title_serif(s, 0.6, 1.05, 12.5, 1.2, "한 주 사용 시나리오", "200% utilization")
add_text(s, 0.6, 2.2, 12.5, 0.5,
         "Mr.AI 도입 후 마케팅 임원의 평균적인 한 주. 모든 단계가 같은 메모리 위에서.",
         font=FONT_KR, size=13, color=TEXT_SOFT)

# Day timeline
days = [
    ("월", "08:00 Daily Briefing", "어제 경쟁사 신상·뉴스 3분 읽기. 이번 주 우선순위 결정."),
    ("화", "콘텐츠 12건 자동 생성", "주제 4 × A/B/C. 페르소나 시뮬 → winner 마킹."),
    ("수", "Market Twin 시뮬", "신규 진출국 후보 검토 · 가격 곡선 도출 (15분)."),
    ("목", "발행 및 디자인 큐레이션", "winner 변형 5 채널 동시 발행. 디자이너는 이미지 review."),
    ("금", "주간 KPI 자동 정리", "이번 주 winner 패턴 학습. 다음 주 spec 자동 업데이트."),
    ("주말", "자동 운영", "크롤 · 브리핑 (월요일분 준비) · 시뮬 큐 자동 진행."),
]
top = 3.3; row_h = 0.55
for i, (d, w, body) in enumerate(days):
    y = top + i * row_h
    # day pill
    add_rect(s, 0.6, y + 0.05, 0.8, 0.4, fill=NAVY)
    add_text(s, 0.6, y + 0.05, 0.8, 0.4, d,
             font=FONT_KR, size=12, color=BG, bold=True,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # work label
    add_text(s, 1.65, y + 0.05, 4.5, 0.4, w,
             font=FONT_KR, size=13, color=NAVY, bold=True,
             anchor=MSO_ANCHOR.MIDDLE)
    # body
    add_text(s, 6.4, y + 0.05, 6.5, 0.4, body,
             font=FONT_KR, size=11, color=TEXT_SOFT,
             anchor=MSO_ANCHOR.MIDDLE)
    # divider
    if i < len(days) - 1:
        add_line(s, 0.6, y + row_h, 12.8, y + row_h,
                 color=LINE, weight=0.4)


# ─────────────────────────────────────────────────────────────────
# Slide 25 — Quantitative results
# ─────────────────────────────────────────────────────────────────
s = new_slide()
page_chrome(s, "20 · 정량 결과", next_pg())
title_serif(s, 0.6, 1.05, 12.5, 1.2, "측정 가능한 결과", "by the numbers")
add_text(s, 0.6, 2.2, 12.5, 0.5,
         "도입 후 평균적인 30일 변화. 정량은 명확, 정성은 더 큼.",
         font=FONT_KR, size=13, color=TEXT_SOFT)

# Big number grid
nums = [
    ("10×", "콘텐츠 생산", "월 8건 → 80건"),
    ("96% ↓", "캠페인 셋업 시간", "1주 → 30분"),
    ("24h", "시장 진출 검증", "기존 6개월"),
    ("3,000", "AI 페르소나 / 시뮬", "Decision+ 티어"),
    ("58% ↓", "연 마케팅 비용", "1.9억 → 0.8억"),
    ("18 / 5", "PostHog 이벤트 / Sentry alert", "정량 운영"),
]
top = 3.3; box_w = 4.0; box_h = 1.85; gap = 0.18
for i, (n, label, sub) in enumerate(nums):
    col = i % 3; row = i // 3
    x = 0.6 + col * (box_w + gap)
    y = top + row * (box_h + 0.18)
    bg_col = ACCENT_SOFT if i == 0 else BG_PURE
    add_rect(s, x, y, box_w, box_h, fill=bg_col, line=LINE, line_w=0.5)
    add_text(s, x + 0.3, y + 0.25, box_w - 0.6, 0.25, label.upper(),
             font=FONT_KR, size=8, color=ACCENT, bold=True, letter_spacing=300)
    # big number serif
    tb = s.shapes.add_textbox(Inches(x + 0.3), Inches(y + 0.55),
                              Inches(box_w - 0.6), Inches(0.85))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]; p.line_spacing = 0.95
    r = p.add_run(); r.text = _clean(n)
    r.font.name = FONT_SERIF
    r.font.size = Pt(40)
    r.font.italic = True
    r.font.color.rgb = NAVY
    add_text(s, x + 0.3, y + 1.45, box_w - 0.6, 0.35, sub,
             font=FONT_KR, size=10.5, color=TEXT_SOFT)


# ─────────────────────────────────────────────────────────────────
# Slide 26 — FSN positioning
# ─────────────────────────────────────────────────────────────────
s = new_slide()
page_chrome(s, "21 · FSN 사업에서의 포지셔닝", next_pg())
title_serif(s, 0.6, 1.05, 12.5, 1.2, "FSN 사업의 한 축", "fsn pillar")
add_text(s, 0.6, 2.2, 12.5, 0.6,
         "FSN 포트폴리오 N 브랜드를 — 같은 인프라로 — 동시에 운영. 브랜드별 격리 · KPI 통합 가시성.",
         font=FONT_KR, size=13, color=TEXT_SOFT)

# Three angles
angles = [
    ("포트폴리오 운영", "N 브랜드 = N 워크스페이스",
     ["각 브랜드의 메모리·KPI·채널 분리",
      "FSN 본사 super admin 은 통합 view",
      "브랜드별 비용 정산 가능"]),
    ("신규 브랜드 onboarding", "Auto-Seed 로 1시간 셋업",
     ["사업계획서 PDF → 메모리 자동",
      "카테고리별 프리셋 자동 적용",
      "1주 안에 콘텐츠 생산 가능"]),
    ("자산 누적 효과", "데이터가 FSN 의 IP",
     ["페르소나 풀 · 톤 학습 · 시뮬 보정",
      "브랜드 종료해도 인사이트 보존",
      "신규 브랜드에 즉시 전이 가능"]),
]
top = 3.3; box_w = 4.0; box_h = 3.2; gap = 0.18
for i, (h, sub, items) in enumerate(angles):
    x = 0.6 + i * (box_w + gap)
    add_rect(s, x, top, box_w, box_h, fill=BG_PURE, line=LINE, line_w=0.5)
    add_pill(s, x + 0.3, top + 0.3, f"0{i+1}",
             fill=ACCENT_SOFT, color=NAVY, size=7, w=0.5)
    add_text(s, x + 0.3, top + 0.75, box_w - 0.6, 0.5, h,
             font=FONT_KR, size=15, color=NAVY, bold=True)
    add_text(s, x + 0.3, top + 1.25, box_w - 0.6, 0.35, sub,
             font=FONT_KR, size=11, color=TEXT_SOFT, italic=True)
    add_bullet_list(s, x + 0.3, top + 1.7, box_w - 0.6, 1.4, items, size=11)


# ─────────────────────────────────────────────────────────────────
# Slide — Next steps (BACK COVER)
# ─────────────────────────────────────────────────────────────────
next_pg()
s = new_slide(bg=NAVY_DEEP)
add_rect(s, 0.6, 1.0, 0.05, 5.3, fill=ACCENT)
add_text(s, 0.95, 1.0, 8, 0.35, "NEXT STEPS",
         font=FONT_KR, size=9, color=ACCENT, bold=True, letter_spacing=400)
add_text(s, 0.95, 1.4, 8, 0.35, "다음 단계",
         font=FONT_KR, size=9, color=RGBColor(0xC9, 0xB7, 0x9A),
         bold=True, letter_spacing=400)

tb = s.shapes.add_textbox(Inches(0.95), Inches(2.0), Inches(11.5), Inches(1.5))
tf = tb.text_frame; tf.word_wrap = True
tf.margin_left = tf.margin_right = 0
p = tf.paragraphs[0]; p.line_spacing = 1.05
r = p.add_run(); r.text = _clean("한 브랜드 · 4주 · ")
r.font.name = FONT_KR; r.font.size = Pt(40); r.font.color.rgb = BG
r2 = p.add_run(); r2.text = "to validate."
r2.font.name = FONT_SERIF; r2.font.size = Pt(40)
r2.font.italic = True; r2.font.color.rgb = RGBColor(0xC9, 0xB7, 0x9A)

# Steps grid
steps_data = [
    ("WEEK 1", "셋업", "워크스페이스 생성 · Auto-Seed · 사업계획서 메모리 시드"),
    ("WEEK 2", "운영 활성화", "크롤·브리핑 가동 · 채널 셋업 · 첫 콘텐츠 12건 생성"),
    ("WEEK 3", "Market Twin 1회", "주요 진출 검토 1건 · 결과 → 콘텐츠 컨텍스트 연결"),
    ("WEEK 4", "KPI 측정 · 결정", "winner 패턴 정리 · ROI 추정 · 정식 도입 결정"),
]
top = 4.3; box_w = 3.0; gap = 0.18
for i, (w, t, body) in enumerate(steps_data):
    x = 0.95 + i * (box_w + gap)
    add_rect(s, x, top, box_w, 2.2,
             fill=RGBColor(0x18, 0x2C, 0x52), line=RGBColor(0x4A, 0x5A, 0x7E))
    add_text(s, x + 0.25, top + 0.25, box_w - 0.5, 0.25, w,
             font=FONT_KR, size=8, color=ACCENT, bold=True, letter_spacing=300)
    add_text(s, x + 0.25, top + 0.6, box_w - 0.5, 0.4, t,
             font=FONT_KR, size=15, color=BG, bold=True)
    add_text(s, x + 0.25, top + 1.15, box_w - 0.5, 1.0, body,
             font=FONT_KR, size=10.5, color=RGBColor(0xC9, 0xB7, 0x9A))


# ─────────────────────────────────────────────────────────────────
# Save
# ─────────────────────────────────────────────────────────────────
import os, sys, time
out = r"C:\Project\Market-Twin\proposals\mrai-core-deck.pptx"
candidates = [out]
# Walk through timestamped fallbacks if the primary or previous
# fallbacks are still locked (PowerPoint open). Each rebuild gets a
# unique name so the user can compare before deciding which to keep.
ts = time.strftime("%H%M%S")
candidates.append(out.replace(".pptx", f"-{ts}.pptx"))
last_err = None
for path in candidates:
    try:
        prs.save(path)
        print(f"saved: {path} ({len(prs.slides)} slides)")
        sys.exit(0)
    except PermissionError as e:
        last_err = e
        print(f"locked: {path}")
print(f"all candidates locked. last error: {last_err}")
sys.exit(1)
