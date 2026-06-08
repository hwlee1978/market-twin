"""
Generate a Toss Payments "결제경로 (빌링)" card-review template (.pptx).

Cover slide is pre-filled with Market Twin merchant info. Slides ②~⑥
each have a header, a one-line capture guide, the URL to capture, and a
large dashed placeholder box where the user pastes the screenshot.

Output: proposals/토스_결제경로_템플릿.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

# ── Palette (Toss blue) ─────────────────────────────────────────
TOSS_BLUE = RGBColor(0x31, 0x82, 0xF6)
TOSS_NAVY = RGBColor(0x19, 0x1F, 0x28)
INK = RGBColor(0x19, 0x1F, 0x28)
TEXT = RGBColor(0x33, 0x3D, 0x4C)
MUTE = RGBColor(0x8B, 0x95, 0xA1)
LINE = RGBColor(0xC9, 0xD3, 0xDE)
BG = RGBColor(0xFF, 0xFF, 0xFF)
BOX_BG = RGBColor(0xF5, 0xF8, 0xFC)

FONT = "맑은 고딕"  # Malgun Gothic — installed on all Windows, renders KR


def set_font(run, name=FONT):
    rPr = run._r.get_or_add_rPr()
    for slot in ("a:latin", "a:ea", "a:cs"):
        el = rPr.find(qn(slot))
        if el is None:
            el = etree.SubElement(rPr, qn(slot))
        el.set("typeface", name)


def add_text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    """runs: list of (text, size, bold, color) — one paragraph each."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, line in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        text, size, bold, color = line
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        set_font(r)
    return tb


def add_rect(slide, x, y, w, h, fill, line_color=None, dashed=False, line_w=1.0):
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    if line_color is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line_color
        sp.line.width = Pt(line_w)
        if dashed:
            ln = sp.line._get_or_add_ln()
            dash = ln.find(qn("a:prstDash"))
            if dash is None:
                dash = etree.SubElement(ln, qn("a:prstDash"))
            dash.set("val", "dash")
    sp.shadow.inherit = False
    return sp


W, H = 13.333, 7.5


def base_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    bg = add_rect(s, -0.1, -0.1, W + 0.2, H + 0.2, BG)
    bg.line.fill.background()
    return s


def capture_slide(prs, step, title, guide, url_label, url, note=None):
    s = base_slide(prs)
    # Header bar
    bar = add_rect(s, 0, 0, W, 1.15, TOSS_BLUE)
    add_text(s, 0.7, 0.0, W - 1.4, 1.15,
             [(f"{step}  {title}", 24, True, RGBColor(0xFF, 0xFF, 0xFF))],
             anchor=MSO_ANCHOR.MIDDLE)
    # Guide line
    add_text(s, 0.7, 1.35, W - 1.4, 0.5,
             [(guide, 14, True, TEXT)])
    # URL line
    add_text(s, 0.7, 1.85, W - 1.4, 0.4,
             [(f"📍 캡처 주소:  {url}", 12, False, TOSS_BLUE)])
    # Placeholder box
    box_y = 2.45
    box_h = 4.3
    add_rect(s, 0.7, box_y, W - 1.4, box_h, BOX_BG, line_color=TOSS_BLUE, dashed=True, line_w=1.5)
    ph = [("▣  화면 캡처를 여기에 붙여넣으세요", 18, True, MUTE),
          ("", 6, False, MUTE),
          ("주소창(도메인) 보이게 · 북마크바 숨김 · PC 시계 함께 · 누락 없이", 12, False, MUTE)]
    if note:
        ph.append(("", 6, False, MUTE))
        ph.append((note, 12, False, TOSS_BLUE))
    add_text(s, 0.7, box_y, W - 1.4, box_h, ph, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    return s


def main():
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)

    # ── Slide 1: Cover (가맹점 정보) ─────────────────────────────
    s = base_slide(prs)
    add_rect(s, 0, 0, 0.35, H, TOSS_BLUE)  # left accent
    add_text(s, 1.0, 1.0, W - 2.0, 1.6,
             [("홈페이지 결제경로 (빌링)", 34, True, INK),
              ("카드사 심사 제출용 · 정기결제(구독)", 16, False, MUTE)])
    # Merchant info box
    bx, by, bw, bh = 1.0, 3.0, 9.5, 3.4
    add_rect(s, bx, by, bw, bh, BOX_BG, line_color=LINE, line_w=1.0)
    add_text(s, bx + 0.5, by + 0.35, bw - 1.0, 0.5,
             [("① 가맹점 정보", 15, True, TOSS_BLUE)])
    info = [
        ("(1) 상호명        :  주식회사 미스터에이아이", 15, False, TEXT),
        ("(2) 사업자번호     :  693-87-03907", 15, False, TEXT),
        ("(3) URL           :  https://markettwin.ai", 15, False, TEXT),
        ("(4) Test ID       :  (테스트 계정 이메일을 입력하세요)", 15, False, MUTE),
        ("(5) Test PW       :  (테스트 계정 비밀번호를 입력하세요)", 15, False, MUTE),
    ]
    add_text(s, bx + 0.5, by + 1.0, bw - 1.0, bh - 1.2, info)
    add_text(s, 1.0, H - 0.7, W - 2.0, 0.4,
             [("2026 · 주식회사 미스터에이아이", 10, False, MUTE)])

    # ── Slides 2~6: capture pages ────────────────────────────────
    capture_slide(
        prs, "②", "하단 정보 캡처",
        "필수 항목: 상호명 · 대표자명 · 사업자등록번호 · 통신판매업신고번호 · 사업장주소 · 유선전화번호",
        "공개 사이트 하단 footer", "markettwin.ai  (하단까지 스크롤 · 로그인 불필요)")

    capture_slide(
        prs, "③", "환불규정 캡처 (무형상품)",
        "정기결제 환불 · 청약철회 제한 · 구독해지 방법 · 자동결제 안내가 모두 보이게",
        "환불정책 페이지 (공개)", "markettwin.ai/refund.html  (로그인 불필요)")

    capture_slide(
        prs, "④", "로그인 / 회원가입 캡처",
        "로그인 폼 + 회원가입 폼 (※ 관리자 '회원가입 공개' 토글 ON 후 캡처)",
        "로그인 / 회원가입", "app.markettwin.ai/login  ·  app.markettwin.ai/signup")

    capture_slide(
        prs, "⑤", "상품 선택 / 구매과정 캡처",
        "로그인 → 결제 메뉴 → 업그레이드 → 플랜 선택 → 주문/결제수단 (필요 시 여러 장)",
        "앱 내 결제 흐름", "app.markettwin.ai  →  결제(/billing) → 플랜(/plans) → /billing/upgrade",
        note="상품 명칭 · 금액 · 결제수단 선택까지 흐름 누락 없이")

    capture_slide(
        prs, "⑥", "카드 결제경로 캡처 (빌링키 발급)",
        "정기결제용 카드 정보 입력창 + 본인인증 화면까지 캡처",
        "Toss 카드등록창 (앱)", "app.markettwin.ai 에서 결제 진행 시 뜨는 Toss 결제창",
        note="⚠ Toss 테스트 키(test_ck_/test_sk_)를 market-twin 프로젝트에 등록 후 가능")

    # ── Slide 7: capture rules reminder ──────────────────────────
    s = base_slide(prs)
    add_rect(s, 0, 0, W, 1.15, TOSS_NAVY)
    add_text(s, 0.7, 0.0, W - 1.4, 1.15,
             [("✓  캡처 공통 규칙 (제출 전 확인)", 22, True, RGBColor(0xFF, 0xFF, 0xFF))],
             anchor=MSO_ANCHOR.MIDDLE)
    rules = [
        ("1.  반드시 PPT 형식으로 제출", 16, True, INK),
        ("2.  모든 캡처에 주소창(도메인)이 보이게 — 북마크바는 숨김", 16, False, TEXT),
        ("3.  모든 캡처에 PC 시계(작업표시줄)가 함께 나오게", 16, False, TEXT),
        ("4.  상품/서비스 명칭·금액·결제창 흐름 누락 없이", 16, False, TEXT),
        ("5.  테스트 결제창 연동으로도 심사 가능 (라이브 승인 전 신청 OK)", 16, False, TEXT),
        ("6.  빌링(정기결제)은 카드 결제창이 아니라 '카드 등록(빌링키 발급)' 창을 캡처", 16, False, TOSS_BLUE),
    ]
    add_text(s, 0.9, 1.8, W - 1.8, 5.0, [(r[0], r[1], r[2], r[3]) for r in rules])
    # spacing between rules
    for p in s.shapes[-1].text_frame.paragraphs:
        p.space_after = Pt(14)

    out = "proposals/토스_결제경로_템플릿_v3.pptx"
    prs.save(out)
    print(f"saved: {out} ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")


if __name__ == "__main__":
    main()
